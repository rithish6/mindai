from __future__ import annotations
import uuid

from fastapi import APIRouter, UploadFile, Depends
from sqlalchemy.orm import Session

from app.models.schemas import Material as MaterialSchema
from app.models.domain import Material
from app.db.database import get_db
from app.core.auth import get_current_user

router = APIRouter()


@router.get("", response_model=list[MaterialSchema])
def list_materials(db: Session = Depends(get_db), current_user: dict = Depends(get_current_user)):
    user_id = current_user.get("uid")
    materials = db.query(Material).filter(Material.user_id == user_id).all()
    return materials


@router.post("/test-upload")
async def test_upload(file: UploadFile):
    return {"filename": file.filename}

@router.get("/debug-token")
def debug_token(auth_header: str = None):
    try:
        from firebase_admin import auth as firebase_auth
        token = auth_header.replace("Bearer ", "") if auth_header else "dummy"
        decoded = firebase_auth.verify_id_token(token)
        return {"status": "success", "decoded": decoded}
    except Exception as e:
        import traceback
        return {"status": "error", "error": str(e), "traceback": traceback.format_exc()}

@router.post("/upload", response_model=MaterialSchema)
async def upload_material(file: UploadFile, db: Session = Depends(get_db), current_user: dict = Depends(get_current_user)):
    user_id = current_user.get("uid")
    filename = file.filename or "uploaded-material"
    filename_lower = filename.lower()
    content_type = file.content_type or ""
    
    file_bytes = await file.read()
    extracted_text = None

    # 1. PDF Documents (.pdf)
    if content_type == "application/pdf" or filename_lower.endswith(".pdf"):
        try:
            import fitz
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            text_pages = [page.get_text() for page in doc if page.get_text().strip()]
            extracted_text = "\n\n".join(text_pages)
            doc.close()
        except Exception as e:
            print(f"Failed to parse PDF via PyMuPDF: {e}")

    # 2. PowerPoint Presentations (.pptx, .ppt)
    elif filename_lower.endswith((".pptx", ".ppt")) or "presentation" in content_type or "powerpoint" in content_type:
        try:
            import io
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            slides_text = []
            for i, slide in enumerate(prs.slides, 1):
                slide_lines = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_lines.append(shape.text.strip())
                if slide_lines:
                    slides_text.append(f"--- Slide {i} ---\n" + "\n".join(slide_lines))
            extracted_text = "\n\n".join(slides_text)
        except Exception as e:
            print(f"Failed to parse PPTX via python-pptx: {e}")

    # 3. Word Documents (.docx, .doc)
    elif filename_lower.endswith((".docx", ".doc")) or "wordprocessing" in content_type or "msword" in content_type:
        try:
            import io
            from docx import Document
            doc = Document(io.BytesIO(file_bytes))
            full_text = [p.text for p in doc.paragraphs if p.text.strip()]
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                    if row_text:
                        full_text.append(row_text)
            extracted_text = "\n".join(full_text)
        except Exception as e:
            print(f"Failed to parse DOCX via python-docx: {e}")

    # 4. Excel Spreadsheets & CSV (.xlsx, .xls, .csv, .tsv)
    elif filename_lower.endswith((".xlsx", ".xls", ".csv", ".tsv")) or "sheet" in content_type or "excel" in content_type or "csv" in content_type:
        try:
            if filename_lower.endswith((".csv", ".tsv")):
                extracted_text = file_bytes.decode("utf-8", errors="ignore")
            else:
                import io
                import openpyxl
                wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
                sheets_text = []
                for sheetname in wb.sheetnames:
                    sheet = wb[sheetname]
                    rows = []
                    for row in sheet.iter_rows(values_only=True):
                        row_cells = [str(c) for c in row if c is not None]
                        if row_cells:
                            rows.append(" | ".join(row_cells))
                    if rows:
                        sheets_text.append(f"--- Sheet: {sheetname} ---\n" + "\n".join(rows))
                extracted_text = "\n\n".join(sheets_text)
        except Exception as e:
            print(f"Failed to parse Spreadsheet: {e}")

    # 5. Images (.png, .jpg, .jpeg, .webp, .bmp) -> Multimodal Vision OCR
    elif content_type.startswith("image/") or filename_lower.endswith((".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif")):
        try:
            from app.services.ai import solve_image_doubt
            extracted_text = solve_image_doubt(
                file_bytes, 
                content_type or "image/png", 
                "Extract and transcribe all educational text, formulas, diagrams, and concepts from this image in detail."
            )
        except Exception as e:
            print(f"Failed image vision OCR: {e}")

    # 6. Audio / Video Media (.mp3, .wav, .m4a, .mp4, .webm, .mov) -> Audio Transcription
    elif content_type.startswith(("audio/", "video/")) or filename_lower.endswith((".mp3", ".wav", ".m4a", ".aac", ".flac", ".mp4", ".mkv", ".mov", ".webm", ".avi")):
        try:
            from app.services.ai import transcribe_media
            extracted_text = transcribe_media(file_bytes, content_type or "audio/mp3")
        except Exception as e:
            print(f"Failed to transcribe media: {e}")

    # 7. Plain Text, Markdown, Code Files (.txt, .md, .json, .csv, .py, .js, etc.)
    if not extracted_text:
        try:
            decoded = file_bytes.decode("utf-8", errors="ignore")
            if decoded.strip() and len(decoded.strip()) > 10:
                extracted_text = decoded
        except Exception as e:
            print(f"Failed UTF-8 decode: {e}")

    # 8. AI Vision / Multimodal Universal Fallback
    if not extracted_text or not extracted_text.strip():
        try:
            # Fallback to AI OCR/text extraction
            from app.services.ai import solve_image_doubt
            extracted_text = solve_image_doubt(
                file_bytes, 
                content_type or "image/png", 
                f"Transcribe and summarize all educational content from this file named '{filename}'."
            )
        except Exception as e:
            print(f"AI fallback extraction failed for {filename}: {e}")

    new_material = Material(
        id=str(uuid.uuid4()),
        user_id=user_id,
        title=filename,
        content=extracted_text if extracted_text else f"Source material content for {filename}",
        material_type=content_type or "application/octet-stream",
        processing_status="processed"
    )
    
    db.add(new_material)
    db.commit()
    db.refresh(new_material)
    
    return new_material
